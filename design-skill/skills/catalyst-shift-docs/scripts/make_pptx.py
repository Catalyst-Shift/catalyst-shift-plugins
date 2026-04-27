"""
Catalyst Shift — HTML deck → PPTX converter
============================================
Reads templates/deck.html (or any HTML with <section data-slide> children),
captures each slide as an image via headless Chrome, and emits a .pptx
with one image per slide at 1280×720 (16:9).

This is the screenshot-based approach — pixel-perfect to the HTML, not
editable as text. Use this when fidelity matters more than post-edit.

Requirements:
    pip install python-pptx playwright
    playwright install chromium

Usage:
    python make_pptx.py path/to/deck.html path/to/output.pptx
"""

import sys
import io
import asyncio
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
    from playwright.async_api import async_playwright
except ImportError as e:
    print(f"Missing dependency: {e}")
    print("Install with: pip install python-pptx playwright && playwright install chromium")
    sys.exit(1)


SLIDE_W = 1280
SLIDE_H = 720


async def capture_slides(html_path: Path) -> list[bytes]:
    """Render HTML and screenshot each <section data-slide> as PNG bytes."""
    url = html_path.resolve().as_uri()
    images: list[bytes] = []
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        ctx = await browser.new_context(viewport={"width": SLIDE_W, "height": SLIDE_H})
        page = await ctx.new_page()
        await page.goto(url, wait_until="networkidle")
        slides = await page.query_selector_all("section[data-slide]")
        if not slides:
            raise RuntimeError("No <section data-slide> elements found.")
        for el in slides:
            png = await el.screenshot(type="png")
            images.append(png)
        await browser.close()
    return images


def build_pptx(images: list[bytes], out_path: Path) -> None:
    prs = Presentation()
    # 16:9 — 13.333" × 7.5"
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for img_bytes in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            io.BytesIO(img_bytes),
            Emu(0), Emu(0),
            width=prs.slide_width, height=prs.slide_height,
        )
    prs.save(str(out_path))


def main() -> None:
    if len(sys.argv) != 3:
        print("Usage: python make_pptx.py <deck.html> <output.pptx>")
        sys.exit(1)
    html_path = Path(sys.argv[1])
    out_path = Path(sys.argv[2])
    if not html_path.exists():
        print(f"HTML not found: {html_path}")
        sys.exit(1)
    images = asyncio.run(capture_slides(html_path))
    build_pptx(images, out_path)
    print(f"Wrote {len(images)} slides → {out_path}")


if __name__ == "__main__":
    main()
